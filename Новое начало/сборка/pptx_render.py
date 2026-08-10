# -*- coding: utf-8 -*-
"""Рендер pptx → HTML → PNG (Chromium) → PDF (Pillow).
Приближение: фигуры-прямоугольники, стрелки, текст с размерами/цветами/поворотами,
таблицы. Шрифт Aptos заменяется системным sans — композиция и цвета сохраняются.
Использование: python3 pptx_render.py file.pptx outdir [--pdf out.pdf]
"""
import sys, os, subprocess, html
from pptx import Presentation
from pptx.util import Emu

SCALE = 160  # px на дюйм (1600x900 на слайд)

def color_of(fill):
    try:
        if fill.type is not None and fill.fore_color and fill.fore_color.type is not None:
            return "#" + str(fill.fore_color.rgb)
    except Exception:
        pass
    return None

def run_html(r):
    t = html.escape(r.text).replace("\n", "<br>")
    size = r.font.size.pt if r.font.size else 12
    color = "#" + str(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else "#262626"
    w = "bold" if r.font.bold else "normal"
    st = "italic" if r.font.italic else "normal"
    fam = r.font.name or "Aptos"
    serif = "serif" in fam.lower()
    ff = "Georgia, serif" if serif else "'Segoe UI', 'DejaVu Sans', Arial, sans-serif"
    sp = ""
    spc = r.font._rPr.get('spc') if r.font._rPr is not None else None
    if spc:
        sp = f"letter-spacing:{int(spc)/100/72*SCALE:.1f}px;"
    return (f'<span style="font-size:{size/72*SCALE:.1f}px;color:{color};'
            f'font-weight:{w};font-style:{st};font-family:{ff};{sp}">{t}</span>')

def para_html(p):
    align = {1: "center", 2: "right", 3: "justify"}.get(p.alignment, "left") if p.alignment else "left"
    ls = p.line_spacing if p.line_spacing else 1.18
    if ls and ls > 3:  # это Pt, не множитель
        ls = 1.18
    runs = "".join(run_html(r) for r in p.runs) or "&nbsp;"
    return f'<div style="text-align:{align};line-height:{ls}">{runs}</div>'

def shape_html(sh):
    out = []
    try:
        x, y = sh.left / 914400 * SCALE, sh.top / 914400 * SCALE
        w, h = sh.width / 914400 * SCALE, sh.height / 914400 * SCALE
    except Exception:
        return out
    rot = getattr(sh, "rotation", 0) or 0
    tf = f"transform:rotate({rot}deg);" if rot else ""
    st = getattr(sh, "shape_type", None)
    stname = str(st) if st is not None else ""
    # заливка фигуры
    if sh.shape_type is not None and not sh.has_text_frame or (sh.has_text_frame and hasattr(sh, "fill")):
        try:
            fill = color_of(sh.fill)
        except Exception:
            fill = None
        if fill:
            if "RIGHT_ARROW" in stname or "ARROW" in stname:
                out.append(f'<div style="position:absolute;left:{x:.0f}px;top:{y:.0f}px;'
                           f'width:{w:.0f}px;height:{h:.0f}px;background:{fill};{tf}'
                           f'clip-path:polygon(0 30%,60% 30%,60% 0,100% 50%,60% 100%,60% 70%,0 70%);"></div>')
            else:
                line = ""
                try:
                    if sh.line and sh.line.color and sh.line.color.type is not None:
                        line = f"border:1px solid #{sh.line.color.rgb};"
                except Exception:
                    pass
                out.append(f'<div style="position:absolute;left:{x:.0f}px;top:{y:.0f}px;'
                           f'width:{w:.0f}px;height:{h:.0f}px;background:{fill};{line}{tf}"></div>')
    # таблица
    if getattr(sh, "has_table", False):
        tbl = sh.table
        rows_html = []
        for row in tbl.rows:
            cells = []
            for cell in row.cells:
                bg = None
                try:
                    bg = color_of(cell.fill)
                except Exception:
                    pass
                content = "".join(para_html(p) for p in cell.text_frame.paragraphs)
                cells.append(f'<td style="background:{bg or "#fff"};padding:4px 10px;'
                             f'border:1px solid #e5e5e5;vertical-align:middle">{content}</td>')
            rows_html.append("<tr>" + "".join(cells) + "</tr>")
        out.append(f'<table style="position:absolute;left:{x:.0f}px;top:{y:.0f}px;'
                   f'width:{w:.0f}px;border-collapse:collapse">' + "".join(rows_html) + "</table>")
        return out
    # текст
    if sh.has_text_frame and sh.text_frame.text.strip():
        va = str(sh.text_frame.vertical_anchor)
        jc = "center" if "MIDDLE" in va else ("flex-end" if "BOTTOM" in va else "flex-start")
        content = "".join(para_html(p) for p in sh.text_frame.paragraphs)
        out.append(f'<div style="position:absolute;left:{x:.0f}px;top:{y:.0f}px;'
                   f'width:{w:.0f}px;height:{h:.0f}px;display:flex;flex-direction:column;'
                   f'justify-content:{jc};overflow:visible;{tf}">{content}</div>')
    return out

def render(pptx_path, outdir, pdf_path=None):
    os.makedirs(outdir, exist_ok=True)
    prs = Presentation(pptx_path)
    W = int(prs.slide_width / 914400 * SCALE)
    H = int(prs.slide_height / 914400 * SCALE)
    pngs = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for sh in slide.shapes:
            try:
                parts += shape_html(sh)
            except Exception as e:
                print(f"  слайд {i}: пропущена фигура ({e})")
        page = (f'<!doctype html><meta charset="utf-8"><body style="margin:0">'
                f'<div style="position:relative;width:{W}px;height:{H}px;'
                f'background:#fff;overflow:hidden;font-family:\'Segoe UI\',Arial,sans-serif">'
                + "".join(parts) + "</div></body>")
        hp = os.path.join(outdir, f"s{i:02d}.html")
        open(hp, "w").write(page)
        png = os.path.join(outdir, f"s{i:02d}.png")
        subprocess.run(["/opt/pw-browsers/chromium", "--headless", "--disable-gpu",
                        "--no-sandbox", f"--window-size={W},{H}",
                        f"--screenshot={os.path.abspath(png)}", f"file://{os.path.abspath(hp)}"],
                       capture_output=True)
        pngs.append(png)
        print("рендер:", png)
    if pdf_path:
        from PIL import Image
        imgs = [Image.open(p).convert("RGB") for p in pngs]
        imgs[0].save(pdf_path, save_all=True, append_images=imgs[1:],
                     resolution=SCALE)
        print("PDF:", pdf_path)

if __name__ == "__main__":
    pptx = sys.argv[1]; outdir = sys.argv[2]
    pdf = None
    if "--pdf" in sys.argv:
        pdf = sys.argv[sys.argv.index("--pdf") + 1]
    render(pptx, outdir, pdf)
