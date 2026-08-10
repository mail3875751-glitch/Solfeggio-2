# -*- coding: utf-8 -*-
"""Внутренняя дека «Время партнёрских конструкций» v2 — те же тексты, что и в
build_deck_internal.py, но каждый слайд 2..11 получил собственную визуальную
форму (дашборды, схемы, диаграммы) вместо повторяющихся карточек 3-в-ряд."""
import sys, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import zipfile, re, os

OUT = sys.argv[1] if len(sys.argv) > 1 else "deck_internal_v2.pptx"

C = dict(dark="1A1A1A", accent="1F3A5F", accent2="2C3E50", body="262626",
         muted="808080", light="F5F5F5", white="FFFFFF", border="D9D9D9",
         warn="C77B30")
F_DISP, F_MAIN, F_SEMI = "Aptos Display", "Aptos", "Aptos SemiBold"
F_SERIF = "Aptos Serif"  # только для цитаты на S7, с фолбэком на Aptos

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

def rgb(h): return RGBColor.from_string(h)
def add_slide(): return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    else: sp.fill.background()
    if line: sp.line.color.rgb = rgb(line); sp.line.width = Pt(0.75)
    else: sp.line.fill.background()
    return sp

def autoshape(s, kind, x, y, w, h, fill=None, line=None, rotation=None):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    else: sp.fill.background()
    if line: sp.line.color.rgb = rgb(line); sp.line.width = Pt(0.75)
    else: sp.line.fill.background()
    if rotation is not None: sp.rotation = rotation
    return sp

def text(s, x, y, w, h, runs, size=13, color=None, bold=False, align="l",
         anchor="t", font=None, caps=False, spacing=None, line_spacing=None,
         italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "c": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    paras = [[(runs, {})]] if isinstance(runs, str) else ([runs] if not (runs and isinstance(runs[0], list)) else runs)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
        if line_spacing: p.line_spacing = line_spacing
        for t, ov in para:
            r = p.add_run()
            r.text = t.upper() if (caps or ov.get("caps")) else t
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.italic = ov.get("italic", italic)
            r.font.name = ov.get("font", font or F_MAIN)
            r.font.color.rgb = rgb(ov.get("color", color or C["body"]))
            sp_val = ov.get("spacing", spacing)
            if sp_val: r.font._rPr.set('spc', str(int(sp_val * 100)))
    return tb

def header(s, label, title, title_size=22):
    text(s, 0.7, 0.42, 8.6, 0.25, label, size=11, color=C["muted"], bold=True,
         font=F_SEMI, caps=True, spacing=3)
    rect(s, 0.7, 0.85, 0.8, 0.045, fill=C["accent"])
    text(s, 0.7, 1.0, 8.6, 0.5, title, size=title_size, color=C["dark"], bold=True, font=F_DISP)

def footer(s, num, source=None):
    if source:
        text(s, 0.7, 4.85, 8.2, 0.25, source, size=10, color=C["muted"], anchor="c")
    text(s, 9.05, 4.85, 0.45, 0.25, str(num), size=10, color=C["muted"], align="r", anchor="c")

def card(s, x, y, w, h, title, body, accent=None, title_size=12.5, body_size=11, num=None):
    accent = accent or C["accent"]
    rect(s, x, y, w, h, fill=C["light"])
    rect(s, x, y, w, 0.055, fill=accent)
    pad, ty = 0.16, y + 0.14
    if num:
        text(s, x + pad, ty, 0.5, 0.32, num, size=18, color=accent, bold=True, font=F_DISP)
        text(s, x + pad + 0.42, ty + 0.02, w - pad * 2 - 0.42, 0.32, title,
             size=title_size, color=C["dark"], bold=True, font=F_SEMI)
        by = ty + 0.36
    else:
        text(s, x + pad, ty, w - pad * 2, 0.3, title, size=title_size,
             color=C["dark"], bold=True, font=F_SEMI)
        by = ty + 0.32
    text(s, x + pad, by, w - pad * 2, y + h - by - 0.12, body, size=body_size,
         color=C["body"], line_spacing=1.12)

def callout(s, label, body, y=None, h=0.5, accent=None):
    accent = accent or C["accent"]
    y = y if y is not None else 4.7 - h
    rect(s, 0.7, y, 8.6, h, fill=C["light"])
    rect(s, 0.7, y, 0.06, h, fill=accent)
    text(s, 0.95, y, 8.2, h, [
        (label + "  ", {"bold": True, "color": accent, "size": 11}),
        (body, {"size": 11.5}),
    ], anchor="c", line_spacing=1.1)

def table(s, x, y, headers, rows, col_w, row_h=0.34, head_h=0.36, fs=10.5, right_cols=()):
    n_r, n_c = len(rows) + 1, len(headers)
    shp = s.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                             Inches(sum(col_w)), Inches(head_h + row_h * len(rows)))
    tbl = shp.table
    tbl.first_row = False; tbl.horz_banding = False
    for j, wcm in enumerate(col_w): tbl.columns[j].width = Inches(wcm)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r): tbl.rows[i].height = Inches(row_h)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(C["accent2"])
        _ct(cell, htxt, fs, C["white"], True, PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(C["light"] if i % 2 else C["white"])
            _ct(cell, val, fs, C["body"], False, PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)
    return shp

def _ct(cell, val, fs, color, bold, align):
    cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for k, chunk in enumerate(val.split("**")):
        if not chunk: continue
        r = p.add_run(); r.text = chunk
        r.font.size = Pt(fs); r.font.bold = bold or (k % 2 == 1)
        r.font.name = F_MAIN; r.font.color.rgb = rgb(color)

# ============ S1 Титул (дословно) ============
s = add_slide()
rect(s, 0, 0, 10, 0.09, fill=C["accent"])
text(s, 0.7, 1.0, 8.6, 0.3, "Внутренний материал · стратегия · август 2026",
     size=12, color=C["muted"], bold=True, font=F_SEMI, caps=True, spacing=3)
rect(s, 0.7, 1.45, 0.8, 0.05, fill=C["accent"])
text(s, 0.7, 1.72, 8.6, 0.85, "Время партнёрских конструкций", size=36,
     color=C["dark"], bold=True, font=F_DISP)
text(s, 0.7, 2.72, 8.0, 0.8,
     "Раньше мы не могли сформулировать, зачем рынку нужны мы. "
     "Теперь рынок сформулировал это за нас — осталось выйти и сказать.",
     size=15, color=C["body"], line_spacing=1.2)
for i in range(5):
    rect(s, 0.7 + i * 0.34, 3.75, 0.22, 0.22, fill=C["border"])
rect(s, 0.7, 4.78, 8.6, 0.012, fill=C["border"])
text(s, 0.7, 4.9, 8.6, 0.25, "Версия 1.0 · август 2026",
     size=10, color=C["muted"])

# ============ S2 «Было → станет» ============
s = add_slide()
header(s, "Контекст", "В эпоху лёгких денег партнёр был не нужен никому")

LX, LY, LW, LH = 0.7, 1.8, 3.4, 2.6
RX, RY, RW, RH = 5.1, 1.6, 4.2, 3.0

rect(s, LX, LY, LW, LH, fill=C["light"])
text(s, LX + 0.2, LY + 0.18, LW - 0.4, 0.5, "2020–2023: деньги за честное слово",
     size=12.5, color=C["muted"], bold=True, font=F_SEMI, line_spacing=1.1)
text(s, LX + 0.2, LY + 0.72, LW - 0.4, LH - 0.9,
     "Девелопер приходил в банк и получал финансирование напрямую — под "
     "проект, под имя, под рост рынка. Бридж под землю был доступен, "
     "льготная ипотека обеспечивала сбыт, маржа прощала ошибки "
     "структурирования.", size=10.5, color=C["muted"], line_spacing=1.2)

# шеврон между панелями, центрирован по вертикальному центру самих панелей
CHV_W, CHV_H = 0.5, 0.4
panel_center_y = (LY + LH / 2 + RY + RH / 2) / 2  # у обеих панелей центр 3.1
chv_cx = (LX + LW + RX) / 2
autoshape(s, MSO_SHAPE.CHEVRON, chv_cx - CHV_W / 2, panel_center_y - CHV_H / 2,
          CHV_W, CHV_H, fill=C["border"])

rect(s, RX, RY, RW, RH, fill=C["light"])
rect(s, RX, RY, RW, 0.055, fill=C["accent"])
text(s, RX + 0.2, RY + 0.2, RW - 0.4, 0.4, "2026: деньги дорогие и недоверчивые",
     size=13, color=C["dark"], bold=True, font=F_SEMI, line_spacing=1.1)
text(s, RX + 0.2, RY + 0.72, RW - 0.4, 1.15,
     "Функция «партнёр-капитал, повышающий надёжность проекта» не имела "
     "спроса: девелоперу партнёры были не нужны, и сформулировать нашу "
     "роль в одном предложении было невозможно. Это не наша ошибка — "
     "это фаза рынка.", size=11, color=C["body"], line_spacing=1.18)
text(s, RX + 0.2, RY + 2.0, RW - 0.4, 0.85,
     "Роль, под которую создавалась структура, ждала своей фазы. Фаза наступила.",
     size=11, color=C["dark"], bold=True, line_spacing=1.15)
footer(s, 2)

# ============ S3 Дашборд рынка ============
s = add_slide()
header(s, "Контекст", "Рынок пересобран: деньги стали дорогими и недоверчивыми")

LBL_X, LBL_W = 0.7, 2.6
SCX, SCW = 3.5, 5.2
OUT_X, OUT_W = 8.75, 0.55
ROW_H, ROW_STEP, ROW0 = 0.5, 0.55, 1.65
MK_W, MK_H = 0.09, 0.28

# каждая строка самодостаточна: подпись значения указана прямо у маркера,
# отдельная легенда «было/сейчас» не нужна
rows_data = [
    dict(label="Ключевая ставка ЦБ", maxv=25, was=21, now=14, single=False,
         vtxt="21% → 14%", outcome="пик пройден"),
    dict(label="Покрытие долга ПФ эскроу", maxv=100, was=88, now=70, single=False,
         vtxt="88% → ~70%", outcome="банки жёстче"),
    dict(label="Объекты с переносом сроков", maxv=100, was=None, now=58, single=True,
         vtxt="сейчас: 58%", outcome="волна реструктуризаций"),
    dict(label="Требование банков к equity", maxv=25, was=None, now=15, single=True,
         vtxt="сейчас: 15%+", outcome="нужен партнёр"),
]

def sc_pos(val, maxv):
    return SCX + (val / maxv) * SCW - MK_W / 2

for i, rd in enumerate(rows_data):
    ry = ROW0 + i * ROW_STEP
    text(s, LBL_X, ry, LBL_W, ROW_H, rd["label"], size=10.5, color=C["body"], anchor="c")
    # значение — сверху; ниже, с зазором, шкала и маркеры, чтобы подпись
    # не перекрывала маркер
    text(s, SCX, ry, SCW, 0.16, rd["vtxt"], size=10, color=C["dark"], bold=True)
    bar_y = ry + 0.26
    rect(s, SCX, bar_y, SCW, 0.16, fill=C["border"])
    mk_y = bar_y - (MK_H - 0.16) / 2
    if not rd["single"]:
        rect(s, sc_pos(rd["was"], rd["maxv"]), mk_y, MK_W, MK_H, fill=C["muted"])
    rect(s, sc_pos(rd["now"], rd["maxv"]), mk_y, MK_W, MK_H, fill=C["accent"])
    text(s, OUT_X, ry, OUT_W, ROW_H, rd["outcome"], size=10, color=C["muted"], anchor="c",
         line_spacing=1.0)

callout(s, "СДВИГ",
        "Девелоперы осваивают навыки банкиров: финансовая инженерия стала ядром бизнеса. "
        "Спрос на переводчика между контурами — максимальный.",
        y=4.08, h=0.62)
footer(s, 3, "Источники: ЦБ РФ, ДОМ.РФ, ЕРЗ.РФ, Коммерсантъ, 2025–2026")

# ============ S4 Навес: диаграмма + KPI ============
s = add_slide()
header(s, "Проблема", "Куплено в 2–3 раза больше, чем рынок может съесть")

BASE_Y = 3.8  # сдвинуто вверх на 0.1" под увеличенный callout
rect(s, 0.7, BASE_Y, 4.4, 0.02, fill=C["border"])

c1_x, c1_w, c1_h = 1.0, 1.3, 0.7
rect(s, c1_x, BASE_Y - c1_h, c1_w, c1_h, fill=C["muted"])
text(s, c1_x, BASE_Y - c1_h, c1_w, c1_h, "1×", size=16, bold=True, color=C["white"],
     align="c", anchor="c", font=F_DISP)

c2_x, c2_w, c2_body_h, c2_top_h = 2.9, 1.3, 1.4, 0.7
rect(s, c2_x, BASE_Y - c2_body_h, c2_w, c2_body_h, fill=C["accent"])
rect(s, c2_x, BASE_Y - c2_body_h - c2_top_h, c2_w, c2_top_h, fill=C["accent2"])
text(s, c2_x, BASE_Y - c2_body_h + 0.15, c2_w, 0.7, "2–3×", size=18, bold=True,
     color=C["white"], align="c", anchor="c", font=F_DISP)

LBL_W = 1.7
c1_cx, c2_cx = c1_x + c1_w / 2, c2_x + c2_w / 2
text(s, c1_cx - LBL_W / 2, BASE_Y + 0.03, LBL_W, 0.24, "Реальный спрос", size=10.5,
     color=C["body"], align="c")
text(s, c2_cx - LBL_W / 2, BASE_Y + 0.03, LBL_W, 0.24, "Купленные площадки", size=10.5,
     color=C["body"], align="c")

KX, KW = 5.5, 3.6
rect(s, KX, 1.5, KW, 1.25, fill=C["light"])
rect(s, KX, 1.5, KW, 0.06, fill=C["warn"])
text(s, KX + 0.18, 1.62, KW - 0.36, 0.55, "52%", size=30, bold=True, color=C["warn"], font=F_DISP)
text(s, KX + 0.18, 2.18, KW - 0.36, 0.5,
     "распроданность домов с вводом 2026 при норме ~70%", size=10.5, color=C["muted"],
     line_spacing=1.15)

rect(s, KX, 2.85, KW, 1.25, fill=C["light"])
rect(s, KX, 2.85, KW, 0.06, fill=C["warn"])
text(s, KX + 0.18, 2.97, KW - 0.36, 0.55, "27%", size=30, bold=True, color=C["warn"], font=F_DISP)
text(s, KX + 0.18, 3.53, KW - 0.36, 0.5,
     "портфеля строится в зоне слабого спроса", size=10.5, color=C["muted"], line_spacing=1.15)

callout(s, "ВОПРОС, КОТОРЫЙ ВСЕ ЗНАЮТ И НИКТО НЕ ПРОИЗНОСИТ",
        "Кто-то должен принять на себя убыток за чрезмерный оптимизм прошлых лет. Кто и в какой форме? "
        "Оценка 2–3× — гипотеза; пилотная экспертиза её проверит.",
        y=4.10, h=0.60, accent=C["warn"])
footer(s, 4, "Оценка масштаба — экспертная; косвенные данные: ДОМ.РФ, ЕРЗ.РФ, 2026")

# ============ S5 Цикл тупика ============
s = add_slide()
header(s, "Проблема", "Убыток есть, но у него нет добровольного хозяина")

B1 = dict(x=0.7, y=1.75, w=2.5, h=0.95)
B2 = dict(x=6.8, y=1.75, w=2.5, h=0.95)
B3 = dict(x=3.75, y=3.35, w=2.5, h=0.95)

def cycle_block(b, title, sub):
    rect(s, b["x"], b["y"], b["w"], b["h"], fill=C["light"])
    rect(s, b["x"], b["y"], b["w"], 0.06, fill=C["accent"])
    text(s, b["x"] + 0.14, b["y"] + 0.14, b["w"] - 0.28, 0.3, title, size=11.5, bold=True,
         color=C["dark"], font=F_SEMI)
    text(s, b["x"] + 0.14, b["y"] + 0.46, b["w"] - 0.28, b["h"] - 0.58, sub, size=10,
         color=C["body"], line_spacing=1.1)

cycle_block(B1, "Девелопер: держит цену", "не может продать ниже «покупка + бридж»")
cycle_block(B2, "Банк: пролонгирует бридж", "забрать залог = резервы по всем аналогичным")
cycle_block(B3, "Оценщик: подтверждает цену", "сопоставимых сделок нет — никто не продаёт")

# стрелка 1→2 (горизонтальная)
a1 = autoshape(s, MSO_SHAPE.RIGHT_ARROW, 3.3, 1.975, 3.4, 0.25, fill=C["muted"])
text(s, 3.3, 1.72, 3.4, 0.2, "цена не падает", size=10, color=C["muted"], align="c")

# центральная плашка — расширена, третий тёмный тон (dark) запрещён, берём accent
PX, PY, PW, PH = 3.9, 2.4, 2.2, 0.8
rect(s, PX, PY, PW, PH, fill=C["accent"])
text(s, PX + 0.1, PY + 0.15, PW - 0.2, 0.25, "EXTEND AND PRETEND", size=10.5, bold=True,
     color=C["white"], align="c", font=F_SEMI, caps=True, spacing=1)
text(s, PX + 0.1, PY + 0.46, PW - 0.2, 0.25, "убыток без хозяина", size=10,
     color=C["border"], align="c")

# стрелка 2→3 (диагональ) — укорочена и сдвинута, чтобы её консервативный
# bbox не задевал плашку и не пересекался с подписью
a2 = autoshape(s, MSO_SHAPE.RIGHT_ARROW, 6.15, 2.9, 0.85, 0.25, fill=C["muted"], rotation=135)
text(s, 7.02, 2.85, 1.4, 0.3, "залог не продаётся", size=10, color=C["muted"], align="l")

# стрелка 3→1 (диагональ) — подпись сдвинута левее консервативного bbox стрелки
a3 = autoshape(s, MSO_SHAPE.RIGHT_ARROW, 2.95, 2.9, 1.0, 0.25, fill=C["muted"], rotation=-145)
text(s, 1.3, 3.00, 1.55, 0.24, "оценка не меняется", size=10, color=C["muted"], align="l")

footer(s, 5)

# ============ S6 Четыре силы ============
s = add_slide()
header(s, "Перелом", "Тупик ломают конструкции, а не ожидание разворота")

small = [
    ("01", "Время", "бриджи имеют сроки; каждая пролонгация дороже предыдущей"),
    ("02", "Регулятор", "ЦБ фиксирует рост реструктуризаций; честность оценок — вопрос времени"),
    ("03", "Первые продажи", "~60 вынужденных сделок M&A уже создают новые ценовые точки"),
]
SB_X, SB_W, SB_H, SB_Y0, SB_STEP = 0.7, 2.9, 0.88, 1.62, 1.0
for i, (num, title, body) in enumerate(small):
    by = SB_Y0 + i * SB_STEP
    rect(s, SB_X, by, SB_W, SB_H, fill=C["light"])
    text(s, SB_X + 0.14, by + 0.12, 0.5, 0.32, num, size=16, bold=True, color=C["muted"], font=F_DISP)
    text(s, SB_X + 0.72, by + 0.14, SB_W - 0.72 - 0.14, 0.32, title, size=11.5, bold=True,
         color=C["dark"], font=F_SEMI)
    text(s, SB_X + 0.14, by + 0.48, SB_W - 0.28, SB_H - 0.58, body, size=10, color=C["body"],
         line_spacing=1.08)

BX, BY, BW, BH = 3.9, 1.62, 5.4, 2.88
rect(s, BX, BY, BW, BH, fill=C["light"])
rect(s, BX, BY, BW, 0.07, fill=C["accent"])
text(s, BX + 0.22, BY + 0.2, 1.2, 0.55, "04", size=30, bold=True, color=C["accent"], font=F_DISP)
text(s, BX + 0.22, BY + 0.78, BW - 0.44, 0.35, "Конструкции", size=15, bold=True,
     color=C["dark"], font=F_SEMI)
text(s, BX + 0.22, BY + 1.18, BW - 0.44, 0.65,
     "механизмы, делающие признание убытка приватным и постепенным — единственный "
     "управляемый путь", size=12, color=C["body"], line_spacing=1.15)
text(s, BX + 0.22, BY + 1.9, BW - 0.44, 0.7,
     "Первые три силы не выбирают выгодоприобретателя. Четвёртая — профессия.",
     size=11, bold=True, color=C["dark"], line_spacing=1.15)
footer(s, 6)

# ============ S7 Цитата отрасли ============
s = add_slide()
header(s, "Перелом", "Отрасль сама назвала тему года")

text(s, 0.7, 1.55, 0.85, 1.0, "«", size=96, bold=True, color=C["border"], font=F_DISP)
text(s, 1.6, 1.9, 7.2, 1.4,
     "От кредитования к совместному бизнесу: эволюция отношений банков и девелоперов",
     size=26, italic=True, color=C["dark"], font=F_SERIF, line_spacing=1.15)
text(s, 1.6, 3.35, 7.2, 0.3,
     "— центральная финансовая конференция форума «Движение», июнь 2026",
     size=12, color=C["muted"])
text(s, 1.6, 4.0, 7.2, 0.4,
     "Центральная тема отрасли — дословно наша заявленная миссия.",
     size=13, bold=True, color=C["accent"])
footer(s, 7)

# ============ S8 Семь конструкций: цвет-секции ============
s = add_slide()
header(s, "Решение", "Семь конструкций: разделить проект, распределить убыток")

def section_plate(x, y, w, h, letter, sq_color, label, fill):
    rect(s, x, y, w, h, fill=fill)
    rect(s, x, y, h, h, fill=sq_color)
    text(s, x, y, h, h, letter, size=13, bold=True, color=C["white"], align="c",
         anchor="c", font=F_DISP)
    text(s, x + h + 0.14, y, w - h - 0.14, h, label, size=11, bold=True, color=C["white"],
         anchor="c", font=F_SEMI, caps=True, spacing=1)

# Заголовок чипа получает фиксированную 2-строчную зону (10pt bold, ~0.32"),
# подзаголовок стартует ниже неё с зазором >=0.05" — не наезжает, даже если
# заголовок реально уместился в одну строку.
CHIP_TOP_PAD, CHIP_TITLE_H, CHIP_GAP, CHIP_BOTTOM_PAD = 0.10, 0.32, 0.05, 0.08

def chip(x, y, w, h, top_color, title, sub):
    rect(s, x, y, w, h, fill=C["light"])
    rect(s, x, y, w, 0.045, fill=top_color)
    text(s, x + 0.12, y + CHIP_TOP_PAD, w - 0.24, CHIP_TITLE_H, title, size=10, bold=True,
         color=C["dark"], font=F_SEMI, line_spacing=1.05)
    sub_y = y + CHIP_TOP_PAD + CHIP_TITLE_H + CHIP_GAP
    sub_h = h - (sub_y - y) - CHIP_BOTTOM_PAD
    text(s, x + 0.12, sub_y, w - 0.24, sub_h, sub, size=10, color=C["muted"],
         line_spacing=1.05)

PA_Y, PA_H = 1.56, 0.30
section_plate(0.7, PA_Y, 8.6, PA_H, "А", C["accent2"], "ВХОД В ПРОЕКТ", C["accent"])
CA_Y, CA_H = PA_Y + PA_H + 0.06, 0.90
chips_a = [
    ("СП с co-investment", "партнёр рискует своими"),
    ("Fee-девелопмент с опционом", "оплата по результату"),
    ("Мезонин с equity-kicker", "долг с конверсией"),
    ("Поэтапный вход по вехам", "переоценка малыми шагами"),
]
cw, gap = 2.05, 0.13
for i, (t, sub) in enumerate(chips_a):
    chip(0.7 + i * (cw + gap), CA_Y, cw, CA_H, C["accent"], t, sub)

PB_Y, PB_H = CA_Y + CA_H + 0.06, 0.30
section_plate(0.7, PB_Y, 8.6, PB_H, "Б", "8F5A20", "РАЗБОР НАВЕСА", C["warn"])
CB_Y, CB_H = PB_Y + PB_H + 0.06, 1.05
chips_b = [
    ("Взнос земли с субординацией", "дисконт через распределение"),
    ("Опцион вместо покупки", "платит тот, кто верит в разворот"),
    ("Конверсия бриджа в долю", "внутригрупповая механика"),
]
cbw, cbgap = 2.77, 0.145
for i, (t, sub) in enumerate(chips_b):
    chip(0.7 + i * (cbw + cbgap), CB_Y, cbw, CB_H, C["warn"], t, sub)

callout(s, "КЛЮЧ",
        "Убыток принимается приватно и постепенно, с сохранением лица всех сторон.",
        y=CB_Y + CB_H + 0.06, h=0.35)
footer(s, 8)

# ============ S9 Bullet graph: покрытие проекта ============
s = add_slide()
header(s, "Решение", "Что получает банк: проект проходит комитет")

text(s, 0.7, 1.54, 8.6, 0.20, "Расчётное покрытие проекта", size=11.5, bold=True,
     color=C["dark"], font=F_SEMI)

SCX2, SCY2, SCW2, SCH2 = 0.7, 2.02, 8.6, 0.30
z1w, z2w, z3w = SCW2 * 0.6, SCW2 * 0.2, SCW2 * 0.2

# подписи зон — НАД шкалой (0.06" зазор до её верхней грани), чтобы не
# перекрываться барами покрытия, которые лежат внутри шкалы
ZLBL_Y, ZLBL_H = 1.80, 0.16
text(s, SCX2, ZLBL_Y, z1w, ZLBL_H, "не проходит", size=10, color=C["muted"], align="c")
text(s, SCX2 + z1w, ZLBL_Y, z2w, ZLBL_H, "зона нормы", size=10, color=C["muted"], align="c")
text(s, SCX2 + z1w + z2w, ZLBL_Y, z3w, ZLBL_H, "комфорт", size=10, color=C["muted"], align="c")

rect(s, SCX2, SCY2, z1w, SCH2, fill="D9D9D9")
rect(s, SCX2 + z1w, SCY2, z2w, SCH2, fill="B1BAC7")
rect(s, SCX2 + z1w + z2w, SCY2, z3w, SCH2, fill="6D7F97")

bar_wo = SCW2 * 0.55
rect(s, SCX2, SCY2 + 0.03, bar_wo, 0.12, fill=C["warn"])
bar_w = SCW2 * 0.78
rect(s, SCX2, SCY2 + 0.17, bar_w, 0.12, fill=C["accent"])

text(s, 0.7, 2.37, 4.2, 0.20, "без партнёрского equity — 55%*", size=10, bold=True, color=C["warn"])
text(s, 4.9, 2.37, 4.4, 0.20, "с партнёрским equity — 78%*", size=10, bold=True, color=C["accent"], align="r")
text(s, 0.7, 2.60, 8.6, 0.16, "* значения иллюстративные", size=10, color=C["muted"])
text(s, 0.7, 2.79, 8.6, 0.16,
     "Два сценария одного проекта на одной шкале — не сложение",
     size=10, color=C["muted"])

card(s, 0.7, 3.02, 4.19, 1.20, "Качество портфеля",
     "Оператор с оплатой за результат снижает риск срыва графика — при 58% "
     "переносов это фактор качества.", body_size=10.5)
card(s, 5.11, 3.02, 4.19, 1.20, "Расчистка навеса",
     "Опционы, субординация, конверсия — без единовременных резервов и "
     "публичных переоценок.", body_size=10.5)

callout(s, "И РОЛЬ",
        "Банк становится партнёром, а не только кредитором, — в русле тренда, "
        "который отрасль объявила главным.", y=4.30, h=0.4)
footer(s, 9, "Значения покрытия на шкале — иллюстративные")

# ============ S10 Схема потоков риска ============
s = add_slide()
header(s, "Наша роль", "Группе нужен приёмник акционерного риска — это мы")

BK = dict(x=0.7, y=1.70, w=2.6, h=0.8)
DC = dict(x=4.6, y=1.70, w=3.0, h=0.8)
rect(s, BK["x"], BK["y"], BK["w"], BK["h"], fill=C["accent"])  # третий тёмный тон запрещён
text(s, BK["x"], BK["y"], BK["w"], BK["h"], "БАНК: кредитный риск", size=11.5, bold=True,
     color=C["white"], align="c", anchor="c", font=F_SEMI, line_spacing=1.05)
rect(s, DC["x"], DC["y"], DC["w"], DC["h"], fill=C["accent"])
text(s, DC["x"], DC["y"], DC["w"], DC["h"], "ПРОФИЛЬНАЯ ДОЧКА — МЫ", size=11.5, bold=True,
     color=C["white"], align="c", anchor="c", font=F_SEMI, line_spacing=1.05)

autoshape(s, MSO_SHAPE.RIGHT_ARROW, 3.4, 1.975, 1.1, 0.25, fill=C["accent"])
text(s, 3.35, 1.62, 1.2, 0.22, "акционерный риск", size=10,
     color=C["accent"], align="c", line_spacing=0.9)
text(s, 3.35, 2.30, 1.2, 0.24, "по утверждённому порядку", size=10, color=C["muted"], align="c")

EB = dict(x=4.6, y=3.00, w=1.85, h=0.55)
PF = dict(x=6.65, y=3.00, w=1.85, h=0.55)
rect(s, EB["x"], EB["y"], EB["w"], EB["h"], fill=C["white"], line=C["border"])
text(s, EB["x"], EB["y"], EB["w"], EB["h"], "Внешний покупатель", size=10, color=C["muted"],
     align="c", anchor="c", line_spacing=1.0)
rect(s, PF["x"], PF["y"], PF["w"], PF["h"], fill=C["light"])
rect(s, PF["x"], PF["y"], PF["w"], 0.05, fill=C["accent"])
text(s, PF["x"], PF["y"] + 0.05, PF["w"], PF["h"] - 0.05, "ПФ-БЛОК", size=11, bold=True,
     color=C["dark"], align="c", anchor="c", font=F_SEMI)

# диагональ "закрыто": банк -> внешний покупатель
bx0, by0 = 2.0, 2.50
bx1, by1 = 4.6, 3.10
dx, dy = bx1 - bx0, by1 - by0
length = math.hypot(dx, dy); ang = math.degrees(math.atan2(dy, dx))
mx, my = (bx0 + bx1) / 2, (by0 + by1) / 2
autoshape(s, MSO_SHAPE.RECTANGLE, mx - length / 2, my - 0.01, length, 0.02,
          fill=C["border"], rotation=ang)
autoshape(s, MSO_SHAPE.RECTANGLE, mx - 0.15, my - 0.025, 0.3, 0.05, fill=C["warn"], rotation=45)
# подпись — ниже консервативного bbox повёрнутой линии и креста (bbox линии
# заканчивается на y≈3.11, креста — на y≈2.93), зазор >=0.05"
text(s, 0.9, 3.16, 3.4, 0.2, "продажа наружу = фиксация убытка — закрыто", size=10,
     color=C["muted"])

# дочка -> ПФ-блок: вертикальная стрелка (без диагональной раздутости bbox),
# в узком "коридоре" x, общем для дочки (4.6-7.6) и ПФ-блока (6.65-8.5)
pf_ax, pf_len, pf_th = 7.1, 0.35, 0.25
pf_ay0, pf_ay1 = DC["y"] + DC["h"] + 0.075, PF["y"] - 0.075  # 2.575 .. 2.925
pf_acy = (pf_ay0 + pf_ay1) / 2
autoshape(s, MSO_SHAPE.RIGHT_ARROW, pf_ax - pf_len / 2, pf_acy - pf_th / 2,
          pf_len, pf_th, fill=C["accent"], rotation=90)
text(s, 7.65, 2.55, 1.65, 0.4, "докапитализация чинит покрытие их кредитов", size=10,
     color=C["muted"], line_spacing=1.0)

text(s, 0.7, 4.30, 8.6, 0.40,
     "Кредитный риск — ваш, акционерный — наш. Мезонин и конверсия бриджа — только "
     "как совместные продукты с ПФ.",
     size=11.5, bold=True, color=C["accent"], align="c", anchor="c", line_spacing=1.05)
footer(s, 10)

# ============ S11 Цифра-герой ============
s = add_slide()
header(s, "Наша роль", "Наше преимущество — взгляд кредитора с обеих сторон")

text(s, 0.7, 1.9, 2.8, 1.35, "20", size=96, bold=True, color=C["accent"], font=F_DISP)
text(s, 0.7, 3.25, 2.8, 0.4, "лет на стыке", size=16, bold=True, color=C["accent"], font=F_SEMI)

theses = [
    "Язык покрытия, ковенант и залогов — родной: мы понимаем и стройку, и кредитный комитет изнутри",
    "Дисциплина «найти, где проект не сойдётся» — готовый due diligence; рынок платит за недоверие",
    "Не всякая площадка должна быть проектом: отличаем переоценённую-но-нужную от ненужной ни за сколько",
]
TX, TW = 3.9, 5.4
TY0, TSTEP = 1.75, 0.95
for i, tstr in enumerate(theses):
    ty = TY0 + i * TSTEP
    if i > 0:
        rect(s, TX, ty - 0.1, TW, 0.012, fill=C["border"])
    text(s, TX, ty, TW, 0.75, tstr, size=12, color=C["body"], line_spacing=1.15, anchor="t")
footer(s, 11)

# ============ S12 Шаги (дословно) ============
s = add_slide()
rect(s, 0, 0, 10, 0.09, fill=C["accent"])
text(s, 0.7, 0.5, 8.6, 0.3, "Предлагаемые шаги", size=12, color=C["muted"],
     bold=True, font=F_SEMI, caps=True, spacing=3)
rect(s, 0.7, 0.95, 0.8, 0.05, fill=C["accent"])
text(s, 0.7, 1.12, 8.6, 0.5, "С чего начать — без риска и без бюджета",
     size=26, color=C["dark"], bold=True, font=F_DISP)
steps = [
    ("Представительство на отраслевых площадках: озвучить повестку конструкций от имени структуры", "сентябрь"),
    ("Пилотная экспертиза 2–3 площадок из навеса: отбор и оценка глазами кредитора", "до 15.10"),
    ("Диалог с блоком проектного финансирования: конструкции как общий инструмент, не конкуренция", "до 31.10"),
    ("Типовая конструкция для кредитного комитета: пакет документов одной сделки-образца", "к декабрю"),
]
y0, STEP = 1.9, 0.56
for i, (t, d) in enumerate(steps):
    yy = y0 + i * STEP
    rect(s, 0.7, yy, 0.36, 0.36, fill=C["accent"])
    text(s, 0.7, yy, 0.36, 0.36, str(i + 1), size=14, color=C["white"],
         bold=True, align="c", anchor="c", font=F_DISP)
    text(s, 1.25, yy - 0.02, 6.8, 0.50, t, size=12.5, color=C["body"], anchor="c",
         line_spacing=1.05)
    text(s, 8.2, yy, 1.1, 0.36, d, size=12, color=C["accent"], bold=True,
         align="r", anchor="c")
text(s, 0.7, 4.20, 8.6, 0.36,
     "При нулевом результате пилота тема закрывается без обязательств и последствий — "
     "риск ограничен временем четырёх шагов.",
     size=11, color=C["muted"], align="c", anchor="c", line_spacing=1.1)
rect(s, 0.7, 4.62, 8.6, 0.012, fill=C["border"])
text(s, 0.7, 4.78, 8.6, 0.3,
     "Ни один шаг не требует аппетита к риску: говорить, отбирать и считать — не инвестировать",
     size=11, color=C["muted"], anchor="c")

prs.save(OUT)

def add_canon_guides(path):
    e8 = lambda inch: round(inch * 576)
    g = ('<p:guideLst>'
         + ''.join(f'<p:guide orient="horz" pos="{e8(y)}"/>' for y in (1.5, 4.7, 4.85, 5.10))
         + ''.join(f'<p:guide pos="{e8(x)}"/>' for x in (0.7, 9.3))
         + '</p:guideLst>')
    tmp = path + '.tmp'
    with zipfile.ZipFile(path) as zin:
        vp = zin.read('ppt/viewProps.xml').decode('utf-8')
        vp = re.sub(r'<p:guideLst>.*?</p:guideLst>|<p:guideLst/>', '', vp, flags=re.S)
        vp = re.sub(r'(</p:cViewPr>)', r'\1' + g, vp, count=1)
        vp = re.sub(r'<p:cSldViewPr(?![^>]*showGuides)', '<p:cSldViewPr showGuides="1"', vp, count=1)
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for n in zin.namelist():
                zout.writestr(n, vp.encode('utf-8') if n == 'ppt/viewProps.xml' else zin.read(n))
    os.replace(tmp, path)

add_canon_guides(OUT)
print("OK:", OUT)

# ================= QA =================
def qa(path):
    from pptx import Presentation as P
    pr = P(path)
    problems = []
    EMU_IN = 914400
    for si, slide in enumerate(pr.slides, start=1):
        if si in (1, 12):
            continue  # титул и финал вне канона
        for shp in slide.shapes:
            top_in = shp.top / EMU_IN if shp.top is not None else 0
            h_in = shp.height / EMU_IN if shp.height is not None else 0
            bottom = top_in + h_in
            if bottom > 5.11:
                problems.append(f"S{si}: shape '{shp.shape_type}' bottom={bottom:.3f} > 5.11")
            if shp.has_text_frame:
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip() and run.font.size is not None:
                            sz = run.font.size.pt
                            if sz < 10:
                                problems.append(f"S{si}: run '{run.text[:30]}' size={sz}pt < 10pt")
    return problems

probs = qa(OUT)
print("QA(a) bounds/font problems:", len(probs))
for p in probs:
    print(" -", p)

EMU_IN = 914400
CANON_SLIDES = set(range(2, 12))  # 2..11


def _text_boxes(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    boxes = []
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            l = shp.left / EMU_IN; t = shp.top / EMU_IN
            w = shp.width / EMU_IN; h = shp.height / EMU_IN
            txt = shp.text_frame.text.strip()[:28] if shp.has_text_frame else ""
            if txt:
                boxes.append((l, t, w, h, txt))
    return boxes


# ---- (б) попарные пересечения независимых текстбоксов, порог 0.03" ----
def overlap_qa(path, slide_nums, thresh=0.03):
    from pptx import Presentation as P
    pr = P(path)
    problems = []
    for si, slide in enumerate(pr.slides, start=1):
        if si not in slide_nums:
            continue
        boxes = _text_boxes(slide)
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                l1, t1, w1, h1, tx1 = boxes[i]
                l2, t2, w2, h2, tx2 = boxes[j]
                ox = min(l1 + w1, l2 + w2) - max(l1, l2)
                oy = min(t1 + h1, t2 + h2) - max(t1, t2)
                if ox > thresh and oy > thresh:
                    problems.append(f"S{si}: '{tx1}' x '{tx2}' overlap ox={ox:.3f} oy={oy:.3f}")
    return problems

oprobs = overlap_qa(OUT, CANON_SLIDES)
print("QA(b) textbox-overlap problems:", len(oprobs))
for p in oprobs:
    print(" -", p)


# ---- (в) фигуры-НЕ-текстбоксы (заливки/стрелки, включая консервативный bbox
#      повёрнутых) пересекающие ЧУЖИЕ независимые текстбоксы, порог 0.05" ----
def shape_text_overlap_qa(path, slide_nums, thresh=0.05, eps=0.02):
    import math as _m
    from pptx import Presentation as P
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    pr = P(path)
    problems = []
    for si, slide in enumerate(pr.slides, start=1):
        if si not in slide_nums:
            continue
        shapes = []
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            if shp.left is None or shp.width is None:
                continue
            l = shp.left / EMU_IN; t = shp.top / EMU_IN
            w = shp.width / EMU_IN; h = shp.height / EMU_IN
            rot = shp.rotation or 0
            cx, cy = l + w / 2, t + h / 2
            if rot % 360 != 0:
                rad = _m.radians(rot)
                hw = (w / 2) * abs(_m.cos(rad)) + (h / 2) * abs(_m.sin(rad))
                hh = (w / 2) * abs(_m.sin(rad)) + (h / 2) * abs(_m.cos(rad))
            else:
                hw, hh = w / 2, h / 2
            bl, bt, bw, bh = cx - hw, cy - hh, hw * 2, hh * 2
            shapes.append((bl, bt, bw, bh))
        boxes = _text_boxes(slide)
        for (sl, st, sw, sh) in shapes:
            for (tl, tt, tw, th, ttxt) in boxes:
                ox = min(sl + sw, tl + tw) - max(sl, tl)
                oy = min(st + sh, tt + th) - max(st, tt)
                if ox <= thresh or oy <= thresh:
                    continue
                # исключение: текст целиком (с допуском eps) лежит внутри фигуры
                # — это подпись на собственной плашке, а не чужой наезд
                contained = (tl >= sl - eps and tt >= st - eps and
                             tl + tw <= sl + sw + eps and tt + th <= st + sh + eps)
                if contained:
                    continue
                problems.append(
                    f"S{si}: shape[{sl:.2f},{st:.2f},{sw:.2f},{sh:.2f}] x '{ttxt}' "
                    f"ox={ox:.3f} oy={oy:.3f}")
    return problems

sprobs = shape_text_overlap_qa(OUT, CANON_SLIDES)
print("QA(c) shape-vs-foreign-text problems:", len(sprobs))
for p in sprobs:
    print(" -", p)

print("QA TOTAL:", len(probs) + len(oprobs) + len(sprobs))
