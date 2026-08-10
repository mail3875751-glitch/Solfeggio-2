# -*- coding: utf-8 -*-
"""Дека «Новое начало» — pptx-strict, профиль neutral, светлый титул."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import zipfile, re, os

OUT = sys.argv[1] if len(sys.argv) > 1 else "deck.pptx"

C = dict(dark="1A1A1A", accent="1F3A5F", accent2="2C3E50", body="262626",
         muted="808080", light="F5F5F5", white="FFFFFF", border="D9D9D9",
         warn="C77B30")

F_DISP, F_MAIN, F_SEMI = "Aptos Display", "Aptos", "Aptos SemiBold"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

def rgb(h): return RGBColor.from_string(h)

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    return sp

def text(s, x, y, w, h, runs, size=13, color=None, bold=False, align="l",
         anchor="t", font=None, caps=False, spacing=None, line_spacing=None,
         space_after=None, wrap=True):
    """runs: str или список (text, {overrides})"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "c": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    if isinstance(runs, str):
        paras = [[(runs, {})]]
    elif runs and isinstance(runs[0], list):
        paras = runs
    else:
        paras = [runs]
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT}[align]
        if line_spacing: p.line_spacing = line_spacing
        if space_after is not None: p.space_after = Pt(space_after)
        for t, ov in para:
            r = p.add_run()
            r.text = t.upper() if (caps or ov.get("caps")) else t
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.name = ov.get("font", font or F_MAIN)
            r.font.color.rgb = rgb(ov.get("color", color or C["body"]))
            if ov.get("italic"): r.font.italic = True
            sp_val = ov.get("spacing", spacing)
            if sp_val:
                r.font._rPr.set('spc', str(int(sp_val * 100)))
    return tb

def header(s, label, title, title_size=22):
    text(s, 0.7, 0.42, 8.6, 0.25, label, size=11, color=C["muted"], bold=True,
         font=F_SEMI, caps=True, spacing=3)
    rect(s, 0.7, 0.85, 0.8, 0.045, fill=C["accent"])
    text(s, 0.7, 1.0, 8.6, 0.5, title, size=title_size, color=C["dark"],
         bold=True, font=F_DISP)

def footer(s, num, source=None):
    if source:
        text(s, 0.7, 4.85, 8.2, 0.25, source, size=10, color=C["muted"], anchor="c")
    text(s, 9.05, 4.85, 0.45, 0.25, str(num), size=10, color=C["muted"],
         align="r", anchor="c")

def card(s, x, y, w, h, title, body, accent=None, title_size=12.5, body_size=11,
         title_color=None, num=None):
    accent = accent or C["accent"]
    rect(s, x, y, w, h, fill=C["light"])
    rect(s, x, y, w, 0.055, fill=accent)
    pad = 0.16
    ty = y + 0.14
    if num:
        text(s, x + pad, ty, 0.5, 0.32, num, size=18, color=accent, bold=True, font=F_DISP)
        text(s, x + pad + 0.42, ty + 0.02, w - pad * 2 - 0.42, 0.32, title, size=title_size,
             color=title_color or C["dark"], bold=True, font=F_SEMI)
        by = ty + 0.36
    else:
        text(s, x + pad, ty, w - pad * 2, 0.3, title, size=title_size,
             color=title_color or C["dark"], bold=True, font=F_SEMI)
        by = ty + 0.32
    text(s, x + pad, by, w - pad * 2, y + h - by - 0.12, body, size=body_size,
         color=C["body"], line_spacing=1.12)

def callout(s, label, body, y=None, h=0.5, accent=None, label_w=None):
    accent = accent or C["accent"]
    y = y if y is not None else 4.7 - h
    rect(s, 0.7, y, 8.6, h, fill=C["light"])
    rect(s, 0.7, y, 0.06, h, fill=accent)
    text(s, 0.95, y, 8.2, h, [
        (label + "  ", {"bold": True, "color": accent, "size": 11}),
        (body, {"size": 11.5}),
    ], anchor="c", line_spacing=1.1)

def kpi(s, x, y, w, h, number, unit, label, num_size=34):
    rect(s, x, y, w, h, fill=C["light"])
    rect(s, x, y, w, 0.055, fill=C["accent"])
    text(s, x + 0.18, y + 0.16, w - 0.36, num_size / 60, [
        (number, {"size": num_size, "bold": True, "color": C["accent"], "font": F_DISP}),
        ("  " + unit, {"size": 15, "bold": True, "color": C["accent"]}),
    ])
    text(s, x + 0.18, y + 0.16 + num_size / 58 + 0.08, w - 0.36, h - (0.16 + num_size / 58 + 0.16),
         label, size=11, color=C["muted"], line_spacing=1.1)

def divider(s, num, title, sub=""):
    rect(s, 0, 0, 10, 5.625, fill=C["accent"])
    text(s, 0.9, 1.55, 8.2, 1.1, num, size=54, color=C["white"], bold=True, font=F_DISP)
    rect(s, 0.95, 2.75, 0.8, 0.045, fill=C["white"])
    text(s, 0.9, 2.95, 8.2, 0.6, title, size=26, color=C["white"], bold=True, font=F_DISP)
    if sub:
        text(s, 0.9, 3.6, 8.2, 0.4, sub, size=13, color="B1BAC7")

def table(s, x, y, headers, rows, col_w, row_h=0.34, head_h=0.36, fs=10.5,
          right_cols=()):
    n_r, n_c = len(rows) + 1, len(headers)
    shp = s.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                             Inches(sum(col_w)), Inches(head_h + row_h * len(rows)))
    tbl = shp.table
    tbl.first_row = False; tbl.horz_banding = False
    for j, wcm in enumerate(col_w):
        tbl.columns[j].width = Inches(wcm)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r):
        tbl.rows[i].height = Inches(row_h)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(C["accent2"])
        _cell_text(cell, htxt, fs, C["white"], True,
                   PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(C["light"] if i % 2 else C["white"])
            _cell_text(cell, val, fs, C["body"], False,
                       PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)
    return shp

def _cell_text(cell, val, fs, color, bold, align):
    cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    parts = val.split("**")
    for k, chunk in enumerate(parts):
        if not chunk: continue
        r = p.add_run(); r.text = chunk
        r.font.size = Pt(fs); r.font.bold = bold or (k % 2 == 1)
        r.font.name = F_MAIN; r.font.color.rgb = rgb(color)

# ================= S1 Титул =================
s = add_slide()
rect(s, 0, 0, 10, 0.09, fill=C["accent"])
text(s, 0.7, 1.0, 8.6, 0.3, "Личная стратегия · август 2026", size=12,
     color=C["muted"], bold=True, font=F_SEMI, caps=True, spacing=3)
rect(s, 0.7, 1.45, 0.8, 0.05, fill=C["accent"])
text(s, 0.7, 1.72, 8.6, 0.85, "Новое начало", size=42, color=C["dark"],
     bold=True, font=F_DISP)
text(s, 0.7, 2.62, 7.6, 0.75,
     "Диагноз ситуации, фаза рынка, траектории и план строительства "
     "собственного имени — на оставшемся runway", size=15, color=C["body"],
     line_spacing=1.2)
for i in range(5):
    rect(s, 0.7 + i * 0.34, 3.75, 0.22, 0.22,
         fill=C["accent"] if i < 2 else C["border"])
rect(s, 0.7, 4.78, 8.6, 0.012, fill=C["border"])
text(s, 0.7, 4.9, 8.6, 0.25, "Рабочая версия 1.0 · для личного пользования",
     size=10, color=C["muted"])

# ================= S2 Резюме =================
s = add_slide()
header(s, "Резюме", "Три вывода, на которых стоит всё остальное")
cw, gap, x0, cy, ch = 2.72, 0.22, 0.7, 1.62, 1.62
card(s, x0, cy, cw, ch, "Эндаумент, не бизнес",
     "Структура живёт на проценты с запаса. Равновесие устраивает того, "
     "у кого право хода. Таймер — уход руководителя.", num="1")
card(s, x0 + cw + gap, cy, cw, ch, "Рынок платит за профиль",
     "В 2026–2027 максимальный спрос — на взгляд кредитора: "
     "реструктуризации, структурирование, вход в землю.", num="2")
card(s, x0 + 2 * (cw + gap), cy, cw, ch, "«4+2» со страховкой",
     "Остаться и строить имя — работает только с датами, метриками "
     "и растяжкой-детектором.", num="3")
kpi(s, 0.7, 3.5, 4.1, 1.06, "2026–2027", "гг.",
    "Окно максимального спроса на профиль. Точка входа — сейчас.")
kpi(s, 5.2, 3.5, 4.1, 1.06, "6–24", "мес.",
    "Строительство видимости до монетизации — во всех траекториях.")
footer(s, 2)

# ================= S3 Divider 01 =================
s = add_slide()
divider(s, "01", "Диагноз", "Структура, равновесие, личная петля")

# ================= S4 Диагноз структуры =================
s = add_slide()
header(s, "Диагноз · структура", "Не пат, а чужое устойчивое равновесие")
cw = 2.72
card(s, 0.7, 1.62, cw, 2.42, "Эндаумент с персоналом",
     "Активная фаза — 2006–2014. Проекты распроданы; запас на депозите "
     "кормит команду. Продукт структуры — комфорт её сотрудников.")
card(s, 0.7 + cw + 0.22, 1.62, cw, 2.42, "Почему банк терпит",
     "Терпимость бесплатна: самофинансирование, прикрытие отношениями "
     "руководителя, статус «опции». Фундамента нет — только отсутствие "
     "повода спросить.")
card(s, 0.7 + 2 * (cw + 0.22), 1.62, cw, 2.42, "Право хода — не моё",
     "Каждый элемент «пата» обслуживает стратегию руководителя "
     "«аккуратно досидеть». Анализ стал генератором оправданий "
     "нерешения.")
callout(s, "КЛЮЧЕВОЙ ВЫВОД",
        "Единственный вопрос в моей зоне контроля: что я строю на оставшемся runway.",
        y=4.22, h=0.48)
footer(s, 4)

# ================= S5 Петля =================
s = add_slide()
header(s, "Диагноз · личный контур", "Одна петля избегания — на трёх масштабах")
rows = [
    ["Неделя", "настроился → не сделал → снова понедельник",
     "задуманное не нужно никому внутри системы"],
    ["8 лет", "ужаснулся (2018) → «не думать» → тот же ужас (2026)",
     "страх оказался точным прогнозом"],
    ["20 лет", "план → расхождение → неудобно перед собой → не смотреть",
     "судья — я двадцатилетней давности"],
]
table(s, 0.7, 1.65, ["Масштаб", "Цикл", "Механизм"], rows,
      [1.15, 4.15, 3.3], row_h=0.5, fs=10.5)
callout(s, "СУТЬ",
        "Катастрофа оказалась не событием, а отсутствием событий: потеря видна "
        "только в сумме, сигнализация не срабатывает на медленную утечку.",
        y=3.95, h=0.62)
footer(s, 5)

# ================= S6 Механизм =================
s = add_slide()
header(s, "Диагноз · личный контур", "Бездействие охраняет старый план")
card(s, 0.7, 1.62, 4.19, 2.3, "Механизм",
     "Пока задуманное не сделано, старый план формально жив. Сделанное "
     "даст ответ реальности — а он может оказаться неудобным: по-старому "
     "не заработает никогда. «Волшебная таблетка» — та же мечта о действии "
     "без действующего.")
card(s, 5.11, 1.62, 4.19, 2.3, "Две ответственности",
     "Перед близкими — несу, и она служит алиби. Перед своей жизнью — "
     "оставлена ~10 лет назад. Вопрос «чего на самом деле ждут близкие — "
     "денег любой ценой или живого меня» — не задан.")
callout(s, "РИСК",
        "Не превращать диагноз в самобичевание: вина — топливо петли. "
        "Подтверждённый прогноз — знание, а не грех.",
        y=4.1, h=0.6, accent=C["warn"])
footer(s, 6)

# ================= S7 Компетенции =================
s = add_slide()
header(s, "Диагноз · компетенции", "В покое компетенции амортизируются быстрее")
card(s, 0.7, 1.62, 4.19, 2.3, "Переворот 1: навык-порок",
     "«Куча аргументов против входа» — болезнь структуры, но "
     "высококлассный личный навык: due diligence глазами банка. "
     "Рынок 2026 платит именно за недоверие к проектам.")
card(s, 5.11, 1.62, 4.19, 2.3, "Переворот 2: жалко терять",
     "Компетенция девелопера умирает не от смены отрасли, а от отсутствия "
     "проектов. Сидеть в бездействующей структуре — режим самой быстрой "
     "потери. «Жалко терять» — аргумент за движение.")
callout(s, "ТРИ НОСИТЕЛЯ ЦЕННОСТИ",
        "Свежесть практики — затухает; имя — конвертировано в бренд руководителя; "
        "соответствие рынку — проверяемо. Лечится только движением.",
        y=4.1, h=0.6)
footer(s, 7)

# ================= S8 Divider 02 =================
s = add_slide()
divider(s, "02", "Рынок", "Фаза цикла, дефицит компетенций, окно")

# ================= S9 Фаза рынка =================
s = add_slide()
header(s, "Рынок · фаза", "Дно пройдено, но экономика проектов ещё ломается")
rows = [
    ["Ключевая ставка ЦБ", "21% → 14%", "10.2024 → 07.2026"],
    ["Рыночная ипотека, рост выдач", "×3,7", "1П 2026 г/г"],
    ["Покрытие долга ПФ эскроу", "88% → ~70%", "2024 → 2026"],
    ["Долг девелоперов по ПФ", "+24%", "прогноз 2026"],
    ["Объекты с переносом сроков", "58%", "2026"],
    ["Сделки M&A (вынужденные)", "~60", "2025"],
]
table(s, 0.7, 1.62, ["Показатель", "Значение", "Период"], rows,
      [4.6, 2.0, 2.0], row_h=0.335, fs=10.5, right_cols=(1,))
callout(s, "ОКНО ПЕРЕРАСПРЕДЕЛЕНИЯ",
        "Кассовый разрыв «продажи вернулись — экономика не починилась» = 12–24 мес. "
        "Слабые распродают проекты сильным.",
        y=4.08, h=0.62)
footer(s, 9, "Источники: ЦБ РФ, ДОМ.РФ, ЕРЗ.РФ, Коммерсантъ, МСФО девелоперов, 2025–2026")

# ================= S10 Дефицит =================
s = add_slide()
header(s, "Рынок · спрос", "Отрасль публично назвала дефицитом мой профиль")
text(s, 0.7, 1.58, 8.6, 0.3,
     "Тема года на форуме «Движение»-2026: «От кредитования к совместному бизнесу»",
     size=12, color=C["muted"])
items = [
    ("Вход в землю", "КРТ, СП, опционы, взнос землёй — без живых денег"),
    ("Банковская финмодель", "покрытие, LLCR, чувствительность — не «бизнес-план»"),
    ("Реструктуризации", "пролонгации, ре-профилирование, waiver"),
    ("Мультибанковость", "организация конкуренции кредиторов"),
    ("Новые инструменты", "ЗПИФ как equity, мезонин, fee-девелопмент"),
]
iw, ih, gx = 2.72, 1.02, 0.22
for i, (t, b) in enumerate(items):
    col, row = i % 3, i // 3
    card(s, 0.7 + col * (iw + gx), 1.98 + row * (ih + 0.16), iw, ih, t, b,
         title_size=11.5, body_size=10)
callout(s, "ВЫВОД",
        "Список читается как описание 20 лет моей работы глазами кредитного комитета.",
        y=4.34, h=0.36)
footer(s, 10)

# ================= S11 Таймеры =================
s = add_slide()
header(s, "Рынок · тайминг", "Три таймера показывают одну точку входа")
card(s, 0.7, 1.62, 2.72, 2.0, "Санаторий",
     "Уход руководителя: разбор структуры, рынок труда в 55+ без имени. "
     "Оценочно 2–5 лет.")
card(s, 3.64, 1.62, 2.72, 2.0, "Кадровый",
     "Имя до монетизации строится 6–24 месяца, и все каналы — "
     "через нетворк.")
card(s, 6.58, 1.62, 2.72, 2.0, "Рыночный",
     "Пик спроса на взгляд кредитора — 2026–2027. К 2028 маржа снова "
     "прикроет неумение структурировать.")
kpi(s, 0.7, 3.82, 8.6, 0.86, "Сейчас", "",
    "Каждый месяц промедления сдвигает монетизацию за горизонт закрытия окна.",
    num_size=26)
footer(s, 11)

# ================= S12 Divider 03 =================
s = add_slide()
divider(s, "03", "Стратегия и план", "Траектории, «4+2», заявка, контрольные точки")

# ================= S13 Кадровая карта =================
s = add_slide()
header(s, "Стратегия · траектории", "Профиль покупается через людей, не по объявлению")
rows = [
    ["CFO / директор по финансированию у девелопера", "300–600 тыс. + бонус",
     "executive search, рекомендации"],
    ["Банк: ПФ и проблемные активы", "150–350 тыс.", "родословная, прямые контакты"],
    ["Интерим на проблемном проекте", "мандаты", "сеть кредитных комитетов"],
    ["CFO крупного девелопера (M&A)", "20–40 млн руб./год", "нужны 2–3 года видимости"],
    ["Портфель независимого директора", "1,5–4 млн руб./год за мандат", "АНД, клуб, преподавание"],
    ["Debt advisory (success fee)", "десятки млн за сделку", "6–18 мес. нулевого дохода"],
]
table(s, 0.7, 1.62, ["Траектория", "Доход (Москва)", "Канал и барьер"], rows,
      [3.7, 2.5, 2.4], row_h=0.335, fs=10)
callout(s, "СКВОЗНОЙ ВЫВОД",
        "Зарплата санатория = финансирование неоплачиваемой фазы. Runway перестаёт "
        "быть анестезией, когда на него покупается конкретное.",
        y=4.08, h=0.62)
footer(s, 13, "Источники: hh.ru, Kontakt InterSearch, НОКС, отраслевые обзоры, авг. 2026")

# ================= S14 Стратегия 4+2 =================
s = add_slide()
header(s, "Стратегия · «4+2»", "«4+2» работает только с датами и метриками")
card(s, 0.7, 1.62, 4.19, 2.3, "Суть",
     "Оставаться, пока санаторий позволяет (4), и строить имя в банковском "
     "контуре и девелоперском сообществе (2). Оплачиваемый runway "
     "финансирует строительство выхода.")
card(s, 5.11, 1.62, 4.19, 2.3, "Главный риск",
     "Единственный вариант, не требующий ничего сегодня. В среде, где "
     "8 лет не выживало ни одно намерение, «4+2» по умолчанию "
     "схлопывается в чистую «4».", accent=C["warn"])
callout(s, "МАНЁВР",
        "Продать руководителю мою видимость как наполнение его экспертного фасада: "
        "ему — легенда без риска, мне — канал в сеть.",
        y=4.1, h=0.6)
footer(s, 14)

# ================= S15 Заявка =================
s = add_slide()
header(s, "Стратегия · заявка", "Имя строится на заявке: «тот, кто ___»")
card(s, 0.7, 1.62, 2.72, 2.2, "Проводник",
     "Проводит девелоперский проект через банковский контур: "
     "структурирование, ПФ, реструктуризация.", num="1")
card(s, 3.64, 1.62, 2.72, 2.2, "Оценщик",
     "Оценивает проекты глазами банка: due diligence, стресс-тест "
     "экономики проекта.", num="2")
card(s, 6.58, 1.62, 2.72, 2.2, "Архитектор партнёрств",
     "Собирает связки банк × девелопер: СП, fee-девелопмент, мезонин, "
     "поэтапный вход.", num="3")
callout(s, "ИНСТРУМЕНТ",
        "«Конструкция-визитка» — не таблетка, а материал для кулуарных разговоров: "
        "каждый показ работает на имя автора.",
        y=4.02, h=0.62)
footer(s, 15)

# ================= S16 Контрольные точки =================
s = add_slide()
header(s, "План · контроль", "Метрики видимости и растяжка-детектор")
rows = [
    ["Выступления / публикации / панели", "2–3", "6+"],
    ["Отношения, проверенные просьбой или помощью", "5", "12"],
    ["Контакты в контуре проектного финансирования", "3", "8"],
    ["Показы «Конструкции-визитки»", "10", "25"],
    ["Входящие обращения за экспертизой", "1–2", "5+"],
]
table(s, 0.7, 1.62, ["Метрика", "6 мес.", "12 мес."], rows,
      [5.6, 1.5, 1.5], row_h=0.335, fs=10.5, right_cols=(1, 2))
callout(s, "РАСТЯЖКА-ДЕТЕКТОР",
        "Нет к 01.03.2027 разговора с руководителем, 3 замеров, 5 показов и одного "
        "публичного появления — включается вариант 3 как есть.",
        y=4.0, h=0.66, accent=C["warn"])
footer(s, 16)

# ================= S17 Финал =================
s = add_slide()
rect(s, 0, 0, 10, 0.09, fill=C["accent"])
text(s, 0.7, 0.5, 8.6, 0.3, "Ближайший месяц", size=12, color=C["muted"],
     bold=True, font=F_SEMI, caps=True, spacing=3)
rect(s, 0.7, 0.95, 0.8, 0.05, fill=C["accent"])
text(s, 0.7, 1.12, 8.6, 0.5, "Следующие шаги", size=28, color=C["dark"],
     bold=True, font=F_DISP)
steps = [
    ("Четыре списка инвентаризации и выбор заявки", "до 22.08"),
    ("«Конструкция-визитка» v0.2 — к первому мероприятию", "до 29.08"),
    ("Два кулуарных мероприятия, по 3 предметных разговора", "до 08.09"),
    ("Замер рынка: 3–4 разговора о цене профиля", "до 30.09"),
    ("Разговор с руководителем о роли лица фасада", "до 30.09"),
]
y0 = 1.85
for i, (t, d) in enumerate(steps):
    yy = y0 + i * 0.52
    rect(s, 0.7, yy, 0.36, 0.36, fill=C["accent"])
    text(s, 0.7, yy, 0.36, 0.36, str(i + 1), size=14, color=C["white"],
         bold=True, align="c", anchor="c", font=F_DISP)
    text(s, 1.25, yy, 6.6, 0.36, t, size=13, color=C["body"], anchor="c")
    text(s, 8.0, yy, 1.3, 0.36, d, size=12, color=C["accent"], bold=True,
         align="r", anchor="c")
rect(s, 0.7, 4.62, 8.6, 0.012, fill=C["border"])
text(s, 0.7, 4.78, 8.6, 0.3,
     "Первый тест стратегии: если разговор с руководителем откладывается — "
     "комбинация уже схлопнулась", size=11, color=C["muted"], anchor="c")

prs.save(OUT)

# Направляющие канона §1.0
def add_canon_guides(path):
    e8 = lambda inch: round(inch * 576)
    g = ('<p:guideLst>'
         + ''.join(f'<p:guide orient="horz" pos="{e8(y)}"/>' for y in (1.5, 4.7, 4.85, 5.10))
         + ''.join(f'<p:guide pos="{e8(x)}"/>' for x in (0.7, 9.3))
         + '</p:guideLst>')
    tmp = path + '.tmp'
    with zipfile.ZipFile(path) as zin:
        vp = zin.read('ppt/viewProps.xml').decode('utf-8')
        vp = re.sub(r'<p:guideLst>.*?</p:guideLst>|<p:guideLst/>', '', vp, flags=re.S)
        vp = re.sub(r'(</p:cViewPr>)', r'\1' + g, vp, count=1)
        vp = re.sub(r'<p:cSldViewPr(?![^>]*showGuides)', '<p:cSldViewPr showGuides="1"', vp, count=1)
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for n in zin.namelist():
                zout.writestr(n, vp.encode('utf-8') if n == 'ppt/viewProps.xml' else zin.read(n))
    os.replace(tmp, path)

add_canon_guides(OUT)
print("OK:", OUT, "slides:", len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst))
