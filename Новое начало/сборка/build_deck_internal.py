# -*- coding: utf-8 -*-
"""Внутренняя дека «Наш момент» — зачем структура нужна рынку сейчас.
Аудитория: руководитель и коллеги. Тон: «мы создавались для этого момента»."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import zipfile, re, os

OUT = sys.argv[1] if len(sys.argv) > 1 else "deck_internal.pptx"

C = dict(dark="1A1A1A", accent="1F3A5F", accent2="2C3E50", body="262626",
         muted="808080", light="F5F5F5", white="FFFFFF", border="D9D9D9",
         warn="C77B30")
F_DISP, F_MAIN, F_SEMI = "Aptos Display", "Aptos", "Aptos SemiBold"

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

def rgb(h): return RGBColor.from_string(h)
def add_slide(): return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    else: sp.fill.background()
    if line: sp.line.color.rgb = rgb(line); sp.line.width = Pt(0.75)
    else: sp.line.fill.background()
    return sp

def text(s, x, y, w, h, runs, size=13, color=None, bold=False, align="l",
         anchor="t", font=None, caps=False, spacing=None, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "c": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    paras = [[(runs, {})]] if isinstance(runs, str) else ([runs] if not (runs and isinstance(runs[0], list)) else runs)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
        if line_spacing: p.line_spacing = line_spacing
        for t, ov in para:
            r = p.add_run()
            r.text = t.upper() if (caps or ov.get("caps")) else t
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.name = ov.get("font", font or F_MAIN)
            r.font.color.rgb = rgb(ov.get("color", color or C["body"]))
            sp_val = ov.get("spacing", spacing)
            if sp_val: r.font._rPr.set('spc', str(int(sp_val * 100)))
    return tb

def header(s, label, title, title_size=22):
    text(s, 0.7, 0.42, 8.6, 0.25, label, size=11, color=C["muted"], bold=True,
         font=F_SEMI, caps=True, spacing=3)
    rect(s, 0.7, 0.85, 0.8, 0.045, fill=C["accent"])
    text(s, 0.7, 1.0, 8.6, 0.5, title, size=title_size, color=C["dark"], bold=True, font=F_DISP)

def footer(s, num, source=None):
    if source:
        text(s, 0.7, 4.85, 8.2, 0.25, source, size=10, color=C["muted"], anchor="c")
    text(s, 9.05, 4.85, 0.45, 0.25, str(num), size=10, color=C["muted"], align="r", anchor="c")

def card(s, x, y, w, h, title, body, accent=None, title_size=12.5, body_size=11, num=None):
    accent = accent or C["accent"]
    rect(s, x, y, w, h, fill=C["light"])
    rect(s, x, y, w, 0.055, fill=accent)
    pad, ty = 0.16, y + 0.14
    if num:
        text(s, x + pad, ty, 0.5, 0.32, num, size=18, color=accent, bold=True, font=F_DISP)
        text(s, x + pad + 0.42, ty + 0.02, w - pad * 2 - 0.42, 0.32, title,
             size=title_size, color=C["dark"], bold=True, font=F_SEMI)
        by = ty + 0.36
    else:
        text(s, x + pad, ty, w - pad * 2, 0.3, title, size=title_size,
             color=C["dark"], bold=True, font=F_SEMI)
        by = ty + 0.32
    text(s, x + pad, by, w - pad * 2, y + h - by - 0.12, body, size=body_size,
         color=C["body"], line_spacing=1.12)

def callout(s, label, body, y=None, h=0.5, accent=None):
    accent = accent or C["accent"]
    y = y if y is not None else 4.7 - h
    rect(s, 0.7, y, 8.6, h, fill=C["light"])
    rect(s, 0.7, y, 0.06, h, fill=accent)
    text(s, 0.95, y, 8.2, h, [
        (label + "  ", {"bold": True, "color": accent, "size": 11}),
        (body, {"size": 11.5}),
    ], anchor="c", line_spacing=1.1)

def table(s, x, y, headers, rows, col_w, row_h=0.34, head_h=0.36, fs=10.5, right_cols=()):
    n_r, n_c = len(rows) + 1, len(headers)
    shp = s.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                             Inches(sum(col_w)), Inches(head_h + row_h * len(rows)))
    tbl = shp.table
    tbl.first_row = False; tbl.horz_banding = False
    for j, wcm in enumerate(col_w): tbl.columns[j].width = Inches(wcm)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r): tbl.rows[i].height = Inches(row_h)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(C["accent2"])
        _ct(cell, htxt, fs, C["white"], True, PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(C["light"] if i % 2 else C["white"])
            _ct(cell, val, fs, C["body"], False, PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)
    return shp

def _ct(cell, val, fs, color, bold, align):
    cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for k, chunk in enumerate(val.split("**")):
        if not chunk: continue
        r = p.add_run(); r.text = chunk
        r.font.size = Pt(fs); r.font.bold = bold or (k % 2 == 1)
        r.font.name = F_MAIN; r.font.color.rgb = rgb(color)

# ============ S1 Титул ============
s = add_slide()
rect(s, 0, 0, 10, 0.09, fill=C["accent"])
text(s, 0.7, 1.0, 8.6, 0.3, "Внутренний материал · стратегия · август 2026",
     size=12, color=C["muted"], bold=True, font=F_SEMI, caps=True, spacing=3)
rect(s, 0.7, 1.45, 0.8, 0.05, fill=C["accent"])
text(s, 0.7, 1.72, 8.6, 0.85, "Время партнёрских конструкций", size=40,
     color=C["dark"], bold=True, font=F_DISP)
text(s, 0.7, 2.62, 8.0, 0.8,
     "Раньше мы не могли сформулировать, зачем рынку нужны мы. "
     "Теперь рынок сформулировал это за нас — осталось выйти и сказать.",
     size=15, color=C["body"], line_spacing=1.2)
for i in range(5):
    rect(s, 0.7 + i * 0.34, 3.75, 0.22, 0.22, fill=C["accent"] if i < 2 else C["border"])
rect(s, 0.7, 4.78, 8.6, 0.012, fill=C["border"])
text(s, 0.7, 4.9, 8.6, 0.25, "Материал для внутреннего обсуждения · v1.0",
     size=10, color=C["muted"])

# ============ S2 Как было ============
s = add_slide()
header(s, "Контекст", "В эпоху лёгких денег партнёр был не нужен никому")
card(s, 0.7, 1.62, 4.19, 2.4, "Как это работало",
     "Девелопер приходил в банк и получал финансирование напрямую — под "
     "проект, под имя, под рост рынка. Бридж под землю был доступен, "
     "льготная ипотека обеспечивала сбыт, маржа прощала ошибки "
     "структурирования.")
card(s, 5.11, 1.62, 4.19, 2.4, "Что это значило для нас",
     "Функция «партнёр-капитал, повышающий надёжность проекта» не имела "
     "спроса: девелоперу партнёры были не нужны, и сформулировать нашу "
     "роль в одном предложении было невозможно. Это не наша ошибка — "
     "это фаза рынка.")
callout(s, "ТЕЗИС",
        "Роль, под которую создавалась структура, ждала своей фазы рынка. Фаза наступила.",
        y=4.22, h=0.48)
footer(s, 2)

# ============ S3 Что изменилось ============
s = add_slide()
header(s, "Контекст", "Рынок пересобран: деньги стали дорогими и недоверчивыми")
rows = [
    ["Ключевая ставка ЦБ", "21% → 14%", "пик пройден, но деньги дороги"],
    ["Покрытие долга ПФ эскроу", "88% → ~70%", "банки ужесточили андеррайтинг"],
    ["Бридж под землю (пик)", "25–26%", "рынок земли в кредит встал"],
    ["Объекты с переносом сроков", "58%", "волна реструктуризаций впереди"],
    ["Сделки M&A за 2025 (вынужденные)", "~60", "перераспределение уже идёт"],
    ["Требование банков к equity", "15%+", "проектам не хватает капитала"],
]
table(s, 0.7, 1.62, ["Показатель", "Значение", "Что это значит"], rows,
      [3.4, 1.6, 3.6], row_h=0.335, fs=10.5, right_cols=(1,))
callout(s, "СДВИГ",
        "Девелоперы осваивают навыки банкиров: финансовая инженерия стала ядром бизнеса. "
        "Спрос на переводчика между контурами — максимальный.",
        y=4.08, h=0.62)
footer(s, 3, "Источники: ЦБ РФ, ДОМ.РФ, ЕРЗ.РФ, Коммерсантъ, NSP, 2025–2026")

# ============ S4 Навес ============
s = add_slide()
header(s, "Проблема", "Земельный навес: куплено кратно больше, чем нужно")
card(s, 0.7, 1.62, 2.72, 2.42, "Масштаб",
     "Потенциал застройки накопленных земельных банков кратно — оценочно "
     "в 2–3 раза — превышает реальный спрос. Куплено в 2020–2023 "
     "на пике, на бриджах.")
card(s, 3.64, 1.62, 2.72, 2.42, "Балансовая иллюзия",
     "Площадки стоят на балансах по ценам рынка, которого больше нет. "
     "Распроданность ввода-2026 — 52% при норме ~70%; 27% портфеля — "
     "в зоне слабого спроса.")
card(s, 6.58, 1.62, 2.72, 2.42, "Вера в разворот",
     "Отрасль ждёт, что рынок «сам съест» навес. Объективно: спрос "
     "меньше потенциального предложения, разворот 2027–2028 съест "
     "лишь часть.")
callout(s, "ВОПРОС, КОТОРЫЙ ВСЕ ЗНАЮТ И НИКТО НЕ ПРОИЗНОСИТ",
        "Кто-то должен принять на себя убыток за чрезмерный оптимизм прошлых лет. Кто и в какой форме?",
        y=4.22, h=0.48, accent=C["warn"])
footer(s, 4, "Оценка масштаба — экспертная; косвенные данные: ДОМ.РФ, ЕРЗ.РФ, 2026")

# ============ S5 Тупик ============
s = add_slide()
header(s, "Проблема", "Убыток есть, но у него нет добровольного хозяина")
card(s, 0.7, 1.62, 2.72, 2.42, "Девелопер",
     "Не может продать ниже «покупка + бридж»: кристаллизация убытка, "
     "ковенанты, дефолт. Держит и пролонгирует.")
card(s, 3.64, 1.62, 2.72, 2.42, "Банк",
     "Не хочет забирать залог: продажа = переоценка всех аналогичных "
     "залогов портфеля и удар резервами. Пролонгировать дешевле.")
card(s, 6.58, 1.62, 2.72, 2.42, "Оценка",
     "«По сопоставимым сделкам», которых нет, потому что никто не "
     "продаёт. Оценка залогов держится на том, что залоги не продаются.")
callout(s, "EXTEND AND PRETEND",
        "Каждый рационален по отдельности — вместе воспроизводят тупик. Признание "
        "убытка публично и разом невозможно ни для кого.",
        y=4.22, h=0.48)
footer(s, 5)

# ============ S6 Что ломает ============
s = add_slide()
header(s, "Перелом", "Тупик ломают конструкции, а не ожидание разворота")
items = [
    ("Время", "бриджи имеют сроки; каждая пролонгация дороже предыдущей", False),
    ("Регулятор", "ЦБ фиксирует рост реструктуризаций; честность залоговых оценок — вопрос времени", False),
    ("Первые продажи", "~60 вынужденных сделок M&A уже создают новые ценовые точки", False),
    ("Конструкции", "механизмы, делающие признание убытка приватным и постепенным — единственный управляемый путь", True),
]
iw, ih = 4.19, 1.02
for i, (t, b, hot) in enumerate(items):
    col, row = i % 2, i // 2
    card(s, 0.7 + col * (iw + 0.22), 1.62 + row * (ih + 0.18), iw, ih, t, b,
         title_size=11.5, body_size=10, accent=C["warn"] if hot else C["accent"])
callout(s, "НАША ТЕРРИТОРИЯ",
        "Первые три силы не выбирают выгодоприобретателя. Четвёртая — профессия: "
        "её делают люди, умеющие считать и договариваться с обеих сторон стола.",
        y=4.12, h=0.58)
footer(s, 6)

# ============ S7 Отраслевой сигнал ============
s = add_slide()
header(s, "Перелом", "Отрасль сама назвала тему года: банки идут в капитал")
text(s, 0.7, 1.7, 8.6, 0.8, [
    [("«От кредитования к совместному бизнесу:", {"size": 20, "bold": True, "color": C["accent"], "font": F_DISP})],
    [("эволюция отношений банков и девелоперов»", {"size": 20, "bold": True, "color": C["accent"], "font": F_DISP})],
], line_spacing=1.1)
text(s, 0.7, 2.62, 8.6, 0.3,
     "— центральная финансовая конференция форума «Движение», июнь 2026",
     size=12, color=C["muted"])
card(s, 0.7, 3.1, 4.19, 1.3, "Что происходит",
     "Банки начали входить в проекты как партнёры и соинвесторы; "
     "fee-девелопмент и СП растут; ЗПИФы недвижимости — ×1,5 за год.",
     body_size=10.5)
card(s, 5.11, 3.1, 4.19, 1.3, "Что это значит для нас",
     "Рынок пришёл в точку, под которую наша структура создавалась: "
     "вход в капитал проектов, повышающий их надёжность для банка.",
     body_size=10.5)
callout(s, "ФОРМУЛА",
        "Центральная тема отрасли — дословно наша заявленная миссия.",
        y=4.52, h=0.4)
footer(s, 7, "Форум «Движение»-2026, июнь 2026")

# ============ S8 Конструкции ============
s = add_slide()
header(s, "Решение", "Семь конструкций: разделить проект, распределить убыток")
text(s, 0.7, 1.56, 8.6, 0.28,
     "Общий принцип: спор о стоимости компетенции и земли переводится из переговоров в механику сделки",
     size=11.5, color=C["muted"])
grp = [
    ("Вход в проект", ["СП с обязательным co-investment",
                       "Fee-девелопмент с опционом на прибыль",
                       "Мезонин с equity-kicker",
                       "Поэтапный вход по вехам"]),
    ("Разбор навеса", ["Взнос земли в СП с субординацией",
                       "Опцион вместо покупки земли",
                       "Конверсия бриджа в долю"]),
]
for gi, (gt, items) in enumerate(grp):
    x = 0.7 + gi * 4.41
    w = 4.19
    rect(s, x, 1.95, w, 2.28, fill=C["light"])
    rect(s, x, 1.95, w, 0.055, fill=C["accent"] if gi == 0 else C["warn"])
    text(s, x + 0.16, 2.08, w - 0.32, 0.28, gt, size=12.5, color=C["dark"],
         bold=True, font=F_SEMI, caps=True, spacing=1)
    rows_txt = [[("—  " + it, {"size": 11})] for it in items]
    text(s, x + 0.16, 2.42, w - 0.32, 1.7, rows_txt, line_spacing=1.35)
callout(s, "КЛЮЧ",
        "Конструкции разбора навеса позволяют принять убыток, не называя его: "
        "приватно, постепенно, с сохранением лица всех сторон.",
        y=4.32, h=0.38)
footer(s, 8)

# ============ S9 Что получает банк ============
s = add_slide()
header(s, "Решение", "Что получает банк: покрытие, портфель, расчистка")
card(s, 0.7, 1.62, 2.72, 2.42, "Покрытие и LLCR",
     "Партнёрское equity в проекте повышает расчётное покрытие — "
     "проходят проекты, которые сегодня получают отказ.")
card(s, 3.64, 1.62, 2.72, 2.42, "Качество портфеля",
     "Оператор с оплатой за результат снижает риск срыва графика — "
     "при 58% переносов это прямой фактор качества.")
card(s, 6.58, 1.62, 2.72, 2.42, "Расчистка навеса",
     "Опционы, субординация, конверсия — расчистка залоговой массы "
     "без единовременных резервов и публичных переоценок.")
callout(s, "И РОЛЬ",
        "Банк становится партнёром, а не только кредитором, — в русле тренда, "
        "который отрасль уже объявила главным.",
        y=4.22, h=0.48)
footer(s, 9)

# ============ S10 Приёмник акционерного риска ============
s = add_slide()
header(s, "Наша роль", "Группе нужен приёмник акционерного риска — это мы")
card(s, 0.7, 1.62, 2.72, 2.42, "Почему не сам банк",
     "Держать equity на балансе банка дорого: нагрузка на капитал, "
     "непрофильный риск. Доли в проектах группы исторически живут "
     "в профильных дочках — это стандартная практика.")
card(s, 3.64, 1.62, 2.72, 2.42, "Почему не внешние",
     "Продажа проблемного актива наружу фиксирует цену и признание "
     "потерь. Внутренняя передача — нет: это «усиление проекта "
     "профильной экспертизой». Рынок навеса закрыт для чужих — "
     "и открыт для нас.")
card(s, 6.58, 1.62, 2.72, 2.42, "Не конкуренция с ПФ",
     "Кредитный риск — их, акционерный — наш: разделение по природе "
     "риска, не по успеху. Наша докапитализация чинит покрытие их же "
     "кредитов — портфель ПФ становится лучше, не меньше.")
callout(s, "РАМКА ЗАХОДА",
        "Мы не просим никого признать провал. Мы принимаем риск, который банку "
        "нельзя держать, — по процедуре, утверждённой сверху, а не по просьбе сбоку.",
        y=4.22, h=0.48)
footer(s, 10)

# ============ S11 Наше преимущество ============
s = add_slide()
header(s, "Наша роль", "Наше преимущество — взгляд кредитора с обеих сторон")
card(s, 0.7, 1.62, 4.19, 1.34, "20 лет на стыке",
     "Мы понимаем и стройку, и кредитный комитет изнутри: язык покрытия, "
     "ковенант, залогов — наш родной.", body_size=10.5)
card(s, 5.11, 1.62, 4.19, 1.34, "Экспертиза отбора",
     "Годы дисциплины «найти, где проект не сойдётся» — это готовый "
     "due diligence. Рынок сейчас платит за недоверие.", body_size=10.5)
card(s, 0.7, 3.12, 4.19, 1.34, "Доверие обеих сторон",
     "Для банка мы — свои по родословной; для девелоперов — отраслевые "
     "эксперты. Позиция модератора переоценки.", body_size=10.5)
card(s, 5.11, 3.12, 4.19, 1.34, "Честная граница",
     "Не всякая площадка должна быть проектом. Отличать «переоценённую, "
     "но нужную» от «ненужной ни за сколько» — наша дисциплина.", body_size=10.5)
footer(s, 10)

# ============ S11 Шаги ============
s = add_slide()
rect(s, 0, 0, 10, 0.09, fill=C["accent"])
text(s, 0.7, 0.5, 8.6, 0.3, "Предлагаемые шаги", size=12, color=C["muted"],
     bold=True, font=F_SEMI, caps=True, spacing=3)
rect(s, 0.7, 0.95, 0.8, 0.05, fill=C["accent"])
text(s, 0.7, 1.12, 8.6, 0.5, "С чего начать — без риска и без бюджета",
     size=26, color=C["dark"], bold=True, font=F_DISP)
steps = [
    ("Представительство на отраслевых площадках: озвучить повестку конструкций от имени структуры", "сентябрь"),
    ("Пилотная экспертиза 2–3 площадок из навеса: отбор и оценка глазами кредитора", "октябрь"),
    ("Диалог с блоком проектного финансирования: конструкции как общий инструмент, не конкуренция", "октябрь"),
    ("Типовая конструкция для кредитного комитета: пакет документов одной сделки-образца", "к декабрю"),
]
y0 = 1.9
for i, (t, d) in enumerate(steps):
    yy = y0 + i * 0.62
    rect(s, 0.7, yy, 0.36, 0.36, fill=C["accent"])
    text(s, 0.7, yy, 0.36, 0.36, str(i + 1), size=14, color=C["white"],
         bold=True, align="c", anchor="c", font=F_DISP)
    text(s, 1.25, yy - 0.02, 6.8, 0.55, t, size=12.5, color=C["body"], anchor="c",
         line_spacing=1.05)
    text(s, 8.2, yy, 1.1, 0.36, d, size=12, color=C["accent"], bold=True,
         align="r", anchor="c")
rect(s, 0.7, 4.62, 8.6, 0.012, fill=C["border"])
text(s, 0.7, 4.78, 8.6, 0.3,
     "Ни один шаг не требует аппетита к риску: говорить, отбирать и считать — не инвестировать",
     size=11, color=C["muted"], anchor="c")

prs.save(OUT)

def add_canon_guides(path):
    e8 = lambda inch: round(inch * 576)
    g = ('<p:guideLst>'
         + ''.join(f'<p:guide orient="horz" pos="{e8(y)}"/>' for y in (1.5, 4.7, 4.85, 5.10))
         + ''.join(f'<p:guide pos="{e8(x)}"/>' for x in (0.7, 9.3))
         + '</p:guideLst>')
    tmp = path + '.tmp'
    with zipfile.ZipFile(path) as zin:
        vp = zin.read('ppt/viewProps.xml').decode('utf-8')
        vp = re.sub(r'<p:guideLst>.*?</p:guideLst>|<p:guideLst/>', '', vp, flags=re.S)
        vp = re.sub(r'(</p:cViewPr>)', r'\1' + g, vp, count=1)
        vp = re.sub(r'<p:cSldViewPr(?![^>]*showGuides)', '<p:cSldViewPr showGuides="1"', vp, count=1)
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for n in zin.namelist():
                zout.writestr(n, vp.encode('utf-8') if n == 'ppt/viewProps.xml' else zin.read(n))
    os.replace(tmp, path)

add_canon_guides(OUT)
print("OK:", OUT)
